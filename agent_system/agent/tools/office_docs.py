"""Инструменты создания офисных документов: Word, Excel, PowerPoint.

Симметрично навыку docparse: тот УМЕЕТ markdown/JSON извлекать ИЗ
docx/xlsx, этот — СОЗДАВАТЬ docx/xlsx/pptx ИЗ markdown/JSON. Результат
docparse можно почти напрямую скормить сюда, если нужно переупаковать
документ или собрать отчёт по данным из нескольких источников.

Никакой LLM внутри — все операции детерминированные (как cad_openscad
или разбор STL): агент передаёт готовый текст/данные, а форматирование
и сохранение файла делает код, а не догадка модели.

Формат ввода для текстовых блоков — упрощённый MARKDOWN (заголовки
`#`/`##`, абзацы, списки `- `/`1. `, таблицы `| a | b |`, **жирный** и
*курсив*). Модели генерируют markdown естественно и без ошибок разметки
чаще, чем произвольный JSON-протокол, поэтому это основной канал ввода;
для таблиц с "неудобным" содержимым (переносы строк, символ `|` в
ячейке) есть отдельные инструменты с явными header/rows в JSON.

python-docx, openpyxl и python-pptx — единственные зависимости,
импортируются ЛЕНИВО (как pymupdf в pdf_pipeline): при отсутствии
библиотеки инструмент вернёт понятную инструкцию по установке.
"""
from __future__ import annotations

import json
import re
from pathlib import Path
from typing import Any

from .base import Tool, ToolError, Workspace


def _require_docx():
    try:
        import docx  # type: ignore
    except ImportError as exc:
        raise ToolError(
            "Создание .docx требует python-docx. Установите: "
            "pip install python-docx"
        ) from exc
    return docx


def _require_openpyxl():
    try:
        import openpyxl  # type: ignore
    except ImportError as exc:
        raise ToolError(
            "Создание .xlsx требует openpyxl. Установите: pip install openpyxl"
        ) from exc
    return openpyxl


def _require_pptx():
    try:
        import pptx  # type: ignore
    except ImportError as exc:
        raise ToolError(
            "Создание .pptx требует python-pptx. Установите: "
            "pip install python-pptx"
        ) from exc
    return pptx


# ======================================================= markdown -> блоки
_INLINE_RE = re.compile(r"(\*\*.+?\*\*|\*.+?\*|`.+?`)")
_TABLE_ROW_RE = re.compile(r"^\s*\|(.+)\|\s*$")
_TABLE_SEP_RE = re.compile(r"^[\s|:\-]+$")
_BULLET_RE = re.compile(r"^(\s*)[-*]\s+(.*)$")
_NUMBERED_RE = re.compile(r"^(\s*)\d+[.)]\s+(.*)$")
_HEADING_RE = re.compile(r"^(#{1,6})\s+(.*)$")


def _split_inline(text: str) -> list[tuple[str, bool, bool, bool]]:
    """Текст -> список (кусок, bold, italic, code)."""
    out: list[tuple[str, bool, bool, bool]] = []
    for tok in _INLINE_RE.split(text):
        if not tok:
            continue
        if tok.startswith("**") and tok.endswith("**") and len(tok) >= 4:
            out.append((tok[2:-2], True, False, False))
        elif tok.startswith("`") and tok.endswith("`") and len(tok) >= 2:
            out.append((tok[1:-1], False, False, True))
        elif tok.startswith("*") and tok.endswith("*") and len(tok) >= 2:
            out.append((tok[1:-1], False, True, False))
        else:
            out.append((tok, False, False, False))
    return out


def parse_markdown_blocks(md: str) -> list[dict[str, Any]]:
    """Упрощённый markdown -> список блоков для docx/pptx.

    Поддержано: заголовки, абзацы (с мягкими переносами строк внутри),
    маркированные/нумерованные списки (с уровнем вложенности по отступу),
    таблицы. Инструкции по конструированию сложных документов (стили,
    картинки) — отдельными инструментами, не через markdown.
    """
    lines = md.replace("\r\n", "\n").split("\n")
    blocks: list[dict[str, Any]] = []
    i, n = 0, len(lines)

    def _level(indent: str) -> int:
        return len(indent) // 2

    while i < n:
        line = lines[i]
        if not line.strip():
            i += 1
            continue

        m = _HEADING_RE.match(line)
        if m:
            blocks.append({"kind": "heading", "level": len(m.group(1)),
                          "text": m.group(2).strip()})
            i += 1
            continue

        if _TABLE_ROW_RE.match(line):
            table_lines = []
            while i < n and _TABLE_ROW_RE.match(lines[i]):
                table_lines.append(lines[i])
                i += 1
            rows = []
            for tl in table_lines:
                inner = _TABLE_ROW_RE.match(tl).group(1)
                if _TABLE_SEP_RE.match(inner):
                    continue  # разделитель |---|---|
                rows.append([c.strip() for c in inner.split("|")])
            if rows:
                blocks.append({"kind": "table", "header": rows[0],
                              "rows": rows[1:]})
            continue

        m = _BULLET_RE.match(line)
        if m:
            items = []
            while i < n and _BULLET_RE.match(lines[i]):
                mm = _BULLET_RE.match(lines[i])
                items.append((_level(mm.group(1)), mm.group(2).strip()))
                i += 1
            blocks.append({"kind": "bullet", "items": items})
            continue

        m = _NUMBERED_RE.match(line)
        if m:
            items = []
            while i < n and _NUMBERED_RE.match(lines[i]):
                mm = _NUMBERED_RE.match(lines[i])
                items.append((_level(mm.group(1)), mm.group(2).strip()))
                i += 1
            blocks.append({"kind": "numbered", "items": items})
            continue

        # обычный абзац: собираем строки до пустой строки/начала другого блока
        para_lines = []
        while i < n and lines[i].strip() and not (
            _HEADING_RE.match(lines[i]) or _TABLE_ROW_RE.match(lines[i])
            or _BULLET_RE.match(lines[i]) or _NUMBERED_RE.match(lines[i])
        ):
            para_lines.append(lines[i].strip())
            i += 1
        blocks.append({"kind": "paragraph", "lines": para_lines})

    return blocks


# ============================================================ docx: запись
def _docx_add_runs(paragraph, text: str, docx_mod) -> None:
    for chunk, bold, italic, code in _split_inline(text):
        run = paragraph.add_run(chunk)
        run.bold = bold or None
        run.italic = italic or None
        if code:
            run.font.name = "Consolas"


def _docx_add_paragraph_lines(doc, lines: list[str], docx_mod) -> None:
    from docx.enum.text import WD_BREAK
    p = doc.add_paragraph()
    for j, ln in enumerate(lines):
        if j > 0:
            p.add_run().add_break(WD_BREAK.LINE)
        _docx_add_runs(p, ln, docx_mod)


def _docx_apply_blocks(doc, blocks: list[dict[str, Any]], docx_mod) -> None:
    for b in blocks:
        kind = b["kind"]
        if kind == "heading":
            level = min(9, max(1, b["level"]))
            doc.add_heading(b["text"], level=level)
        elif kind == "paragraph":
            _docx_add_paragraph_lines(doc, b["lines"], docx_mod)
        elif kind == "bullet":
            for lvl, text in b["items"]:
                style = "List Bullet" if lvl == 0 else f"List Bullet {min(lvl + 1, 3)}"
                try:
                    p = doc.add_paragraph(style=style)
                except KeyError:
                    p = doc.add_paragraph(style="List Bullet")
                _docx_add_runs(p, text, docx_mod)
        elif kind == "numbered":
            for lvl, text in b["items"]:
                style = "List Number" if lvl == 0 else f"List Number {min(lvl + 1, 3)}"
                try:
                    p = doc.add_paragraph(style=style)
                except KeyError:
                    p = doc.add_paragraph(style="List Number")
                _docx_add_runs(p, text, docx_mod)
        elif kind == "table":
            _docx_add_table_rows(doc, b["header"], b["rows"], "Table Grid", docx_mod)


def _docx_add_table_rows(doc, header: list[str], rows: list[list[str]],
                         style: str, docx_mod) -> None:
    cols = len(header) or (len(rows[0]) if rows else 1)
    table = doc.add_table(rows=1, cols=cols)
    try:
        table.style = style
    except KeyError as exc:
        raise ToolError(
            f"Стиль таблицы {style!r} неизвестен шаблону Word. Оставьте "
            "по умолчанию 'Table Grid' или используйте валидный стиль."
        ) from exc
    for c, val in enumerate(header):
        table.rows[0].cells[c].text = str(val)
    for row in rows:
        cells = table.add_row().cells
        for c in range(cols):
            cells[c].text = str(row[c]) if c < len(row) else ""


# ============================================================ pptx: запись
def _pptx_layout(prs, name: str):
    names = {
        "title": 0, "title_content": 1, "section_header": 2,
        "two_content": 3, "title_only": 5, "blank": 6,
    }
    idx = names.get(name)
    if idx is None or idx >= len(prs.slide_layouts):
        raise ToolError(
            f"Неизвестная раскладка {name!r}. Доступны: {', '.join(names)}"
        )
    return prs.slide_layouts[idx]


def _pptx_fill_body(text_frame, blocks: list[dict[str, Any]]) -> None:
    first = True
    for b in blocks:
        if b["kind"] == "paragraph":
            text = " ".join(b["lines"])
            p = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
            p.text = text
            first = False
        elif b["kind"] in ("bullet", "numbered"):
            for lvl, text in b["items"]:
                p = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
                p.text = text
                p.level = min(lvl, 4)
                first = False
        elif b["kind"] == "heading":
            p = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
            p.text = b["text"]
            first = False
        # таблицы на слайде — через pptx_add_table_slide, здесь пропускаем


# ================================================================== build
def build(ws: Workspace) -> list[Tool]:

    # ------------------------------------------------------------- docx
    def docx_create(path: str, title: str = "", markdown: str = "") -> str:
        docx_mod = _require_docx()
        p = ws.resolve(path)
        p.parent.mkdir(parents=True, exist_ok=True)
        doc = docx_mod.Document()
        if title.strip():
            doc.add_heading(title.strip(), level=0)
        if markdown.strip():
            _docx_apply_blocks(doc, parse_markdown_blocks(markdown), docx_mod)
        doc.save(p)
        return f"Создан {ws.relative(p)} ({len(doc.paragraphs)} абзацев)"

    def docx_append(path: str, markdown: str) -> str:
        if not markdown.strip():
            raise ToolError("markdown пуст — нечего добавлять")
        docx_mod = _require_docx()
        p = ws.resolve(path)
        p.parent.mkdir(parents=True, exist_ok=True)
        doc = docx_mod.Document(str(p)) if p.exists() else docx_mod.Document()
        before = len(doc.paragraphs)
        _docx_apply_blocks(doc, parse_markdown_blocks(markdown), docx_mod)
        doc.save(p)
        return (f"{'Дополнен' if p.exists() else 'Создан'} {ws.relative(p)}: "
               f"+{len(doc.paragraphs) - before} абзацев")

    def docx_add_table(path: str, header_json: str, rows_json: str,
                       style: str = "Table Grid") -> str:
        docx_mod = _require_docx()
        try:
            header = json.loads(header_json)
            rows = json.loads(rows_json) if rows_json.strip() else []
        except json.JSONDecodeError as exc:
            raise ToolError(f"header_json/rows_json должны быть JSON: {exc}") from exc
        if not isinstance(header, list) or not isinstance(rows, list):
            raise ToolError("header_json — список столбцов, rows_json — список строк")
        p = ws.resolve(path)
        p.parent.mkdir(parents=True, exist_ok=True)
        doc = docx_mod.Document(str(p)) if p.exists() else docx_mod.Document()
        _docx_add_table_rows(doc, header, rows, style, docx_mod)
        doc.save(p)
        return f"Таблица добавлена в {ws.relative(p)}: {len(rows)} строк данных"

    def docx_add_image(path: str, image_path: str, width_cm: float = 0) -> str:
        docx_mod = _require_docx()
        from docx.shared import Cm
        p = ws.resolve(path)
        img = ws.resolve(image_path)
        if not img.exists():
            raise ToolError(f"Изображение {image_path!r} не найдено")
        p.parent.mkdir(parents=True, exist_ok=True)
        doc = docx_mod.Document(str(p)) if p.exists() else docx_mod.Document()
        if width_cm and width_cm > 0:
            doc.add_picture(str(img), width=Cm(width_cm))
        else:
            doc.add_picture(str(img))
        doc.save(p)
        return f"Изображение {ws.relative(img)} добавлено в {ws.relative(p)}"

    # ------------------------------------------------------------- xlsx
    def xlsx_create(path: str, sheets_json: str) -> str:
        openpyxl = _require_openpyxl()
        try:
            sheets = json.loads(sheets_json)
        except json.JSONDecodeError as exc:
            raise ToolError(f"sheets_json должен быть JSON: {exc}") from exc
        if not isinstance(sheets, dict) or not sheets:
            raise ToolError(
                'sheets_json должен быть объектом {"Лист": [[строка...],...]}'
            )
        p = ws.resolve(path)
        p.parent.mkdir(parents=True, exist_ok=True)
        wb = openpyxl.Workbook()
        wb.remove(wb.active)
        for name, rows in sheets.items():
            if not isinstance(rows, list):
                raise ToolError(f"Лист {name!r}: строки должны быть списком списков")
            sh = wb.create_sheet(str(name)[:31])
            for row in rows:
                sh.append(row if isinstance(row, list) else [row])
        wb.save(p)
        return f"Создан {ws.relative(p)}: {len(sheets)} лист(ов)"

    def xlsx_write_rows(path: str, sheet: str, rows_json: str,
                        start_row: int = 0, header_bold: bool = False) -> str:
        openpyxl = _require_openpyxl()
        from openpyxl.styles import Font
        try:
            rows = json.loads(rows_json)
        except json.JSONDecodeError as exc:
            raise ToolError(f"rows_json должен быть JSON-списком списков: {exc}") from exc
        if not isinstance(rows, list) or (rows and not isinstance(rows[0], list)):
            raise ToolError("rows_json должен быть списком списков (строк таблицы)")
        p = ws.resolve(path)
        p.parent.mkdir(parents=True, exist_ok=True)
        wb = openpyxl.load_workbook(p) if p.exists() else openpyxl.Workbook()
        if not p.exists() and wb.active.title == "Sheet" and sheet != "Sheet":
            wb.active.title = sheet
        sh = wb[sheet] if sheet in wb.sheetnames else wb.create_sheet(sheet)
        r0 = start_row if start_row > 0 else sh.max_row + 1 if sh.max_row > 1 or \
            any(c.value is not None for c in sh[1]) else 1
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = sh.cell(row=r0 + i, column=j + 1, value=val)
                if header_bold and i == 0:
                    cell.font = Font(bold=True)
        wb.save(p)
        return (f"Записано {len(rows)} строк в {ws.relative(p)}::{sheet} "
               f"начиная со строки {r0}")

    def xlsx_add_chart(path: str, sheet: str, chart_type: str, data_range: str,
                       categories_range: str = "", title: str = "",
                       anchor: str = "E2") -> str:
        openpyxl = _require_openpyxl()
        from openpyxl.chart import BarChart, LineChart, PieChart, Reference
        from openpyxl.utils.cell import range_boundaries
        p = ws.resolve(path)
        if not p.exists():
            raise ToolError(f"Файл {path!r} не найден — сначала создайте его")
        wb = openpyxl.load_workbook(p)
        if sheet not in wb.sheetnames:
            raise ToolError(f"Лист {sheet!r} не найден. Есть: {', '.join(wb.sheetnames)}")
        sh = wb[sheet]

        charts = {"bar": BarChart, "line": LineChart, "pie": PieChart}
        cls = charts.get(chart_type)
        if cls is None:
            raise ToolError(f"Тип диаграммы {chart_type!r}. Доступны: {', '.join(charts)}")
        chart = cls()
        if title:
            chart.title = title

        try:
            min_c, min_r, max_c, max_r = range_boundaries(data_range)
        except ValueError as exc:
            raise ToolError(f"data_range {data_range!r} невалиден: {exc}") from exc
        data = Reference(sh, min_col=min_c, min_row=min_r, max_col=max_c, max_row=max_r)
        chart.add_data(data, titles_from_data=True)
        if categories_range.strip():
            try:
                cmin_c, cmin_r, cmax_c, cmax_r = range_boundaries(categories_range)
            except ValueError as exc:
                raise ToolError(
                    f"categories_range {categories_range!r} невалиден: {exc}") from exc
            cats = Reference(sh, min_col=cmin_c, min_row=cmin_r,
                             max_col=cmax_c, max_row=cmax_r)
            chart.set_categories(cats)
        sh.add_chart(chart, anchor)
        wb.save(p)
        return f"Диаграмма {chart_type} добавлена на {ws.relative(p)}::{sheet}@{anchor}"

    # ------------------------------------------------------------- pptx
    def pptx_create(path: str, title: str = "", subtitle: str = "") -> str:
        pptx_mod = _require_pptx()
        p = ws.resolve(path)
        p.parent.mkdir(parents=True, exist_ok=True)
        prs = pptx_mod.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if title.strip():
            slide.shapes.title.text = title.strip()
        if subtitle.strip() and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle.strip()
        prs.save(p)
        return f"Создана презентация {ws.relative(p)}: 1 слайд (титульный)"

    def pptx_add_slide(path: str, title: str = "", markdown_body: str = "",
                       notes: str = "", layout: str = "title_content") -> str:
        pptx_mod = _require_pptx()
        p = ws.resolve(path)
        p.parent.mkdir(parents=True, exist_ok=True)
        prs = pptx_mod.Presentation(str(p)) if p.exists() else pptx_mod.Presentation()
        slide = prs.slides.add_slide(_pptx_layout(prs, layout))
        if title.strip() and slide.shapes.title is not None:
            slide.shapes.title.text = title.strip()
        if markdown_body.strip():
            body_ph = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx != 0:  # не заголовок
                    body_ph = shape
                    break
            if body_ph is None:
                raise ToolError(
                    f"Раскладка {layout!r} не содержит места для текста. "
                    "Используйте layout='title_content'."
                )
            _pptx_fill_body(body_ph.text_frame, parse_markdown_blocks(markdown_body))
        if notes.strip():
            slide.notes_slide.notes_text_frame.text = notes.strip()
        prs.save(p)
        return (f"Слайд добавлен в {ws.relative(p)}: "
               f"всего {len(prs.slides)} слайд(ов)")



    def pptx_add_table_slide(path: str, title: str, header_json: str,
                             rows_json: str) -> str:
        pptx_mod = _require_pptx()
        from pptx.util import Inches
        try:
            header = json.loads(header_json)
            rows = json.loads(rows_json) if rows_json.strip() else []
        except json.JSONDecodeError as exc:
            raise ToolError(f"header_json/rows_json должны быть JSON: {exc}") from exc
        if not isinstance(header, list):
            raise ToolError("header_json должен быть списком названий столбцов")
        p = ws.resolve(path)
        p.parent.mkdir(parents=True, exist_ok=True)
        prs = pptx_mod.Presentation(str(p)) if p.exists() else pptx_mod.Presentation()
        slide = prs.slides.add_slide(_pptx_layout(prs, "title_only"))
        if title.strip() and slide.shapes.title is not None:
            slide.shapes.title.text = title.strip()
        n_rows, n_cols = len(rows) + 1, len(header)
        table_shape = slide.shapes.add_table(
            n_rows, n_cols, Inches(0.5), Inches(1.6), Inches(9), Inches(0.4 * n_rows))
        table = table_shape.table
        for c, val in enumerate(header):
            table.cell(0, c).text = str(val)
        for r, row in enumerate(rows, start=1):
            for c in range(n_cols):
                table.cell(r, c).text = str(row[c]) if c < len(row) else ""
        prs.save(p)
        return f"Слайд с таблицей ({len(rows)} строк) добавлен в {ws.relative(p)}"

    def pptx_add_image_slide(path: str, image_path: str, title: str = "",
                             notes: str = "") -> str:
        pptx_mod = _require_pptx()
        from pptx.util import Inches
        p = ws.resolve(path)
        img = ws.resolve(image_path)
        if not img.exists():
            raise ToolError(f"Изображение {image_path!r} не найдено")
        p.parent.mkdir(parents=True, exist_ok=True)
        prs = pptx_mod.Presentation(str(p)) if p.exists() else pptx_mod.Presentation()
        layout_name = "title_only" if title.strip() else "blank"
        slide = prs.slides.add_slide(_pptx_layout(prs, layout_name))
        if title.strip() and slide.shapes.title is not None:
            slide.shapes.title.text = title.strip()
        top = Inches(1.6) if title.strip() else Inches(0.5)
        slide.shapes.add_picture(str(img), Inches(0.5), top, width=Inches(9))
        if notes.strip():
            slide.notes_slide.notes_text_frame.text = notes.strip()
        prs.save(p)
        return f"Слайд с изображением {ws.relative(img)} добавлен в {ws.relative(p)}"

    return [
        Tool("docx_create",
             "Создать новый документ Word. Тело — упрощённый markdown: "
             "заголовки '#'/'##', абзацы, списки '- '/'1. ', таблицы "
             "'| a | b |', **жирный**/*курсив*. Перезаписывает файл, если он "
             "уже существует — для дополнения используйте docx_append.",
             {"type": "object",
              "properties": {
                  "path": {"type": "string"},
                  "title": {"type": "string", "description": "Заголовок титульного уровня"},
                  "markdown": {"type": "string"}},
              "required": ["path"]},
             docx_create),
        Tool("docx_append",
             "Добавить содержимое (markdown) в конец существующего Word-"
             "документа; если файла нет — создаёт новый. Так документ можно "
             "собирать по частям за несколько вызовов.",
             {"type": "object",
              "properties": {
                  "path": {"type": "string"},
                  "markdown": {"type": "string"}},
              "required": ["path", "markdown"]},
             docx_append),
        Tool("docx_add_table",
             "Добавить таблицу в конец документа Word из явных JSON-списков "
             "(header_json — столбцы, rows_json — строки данных). Используйте "
             "вместо markdown-таблицы, если в ячейках есть символ '|' или "
             "переносы строк.",
             {"type": "object",
              "properties": {
                  "path": {"type": "string"},
                  "header_json": {"type": "string", "description": '["Поз","Наименование"]'},
                  "rows_json": {"type": "string", "description": '[["1","Корпус"],...]'},
                  "style": {"type": "string", "description": "Стиль таблицы Word"}},
              "required": ["path", "header_json"]},
             docx_add_table),
        Tool("docx_add_image",
             "Добавить изображение (из рабочей папки) в конец документа Word.",
             {"type": "object",
              "properties": {
                  "path": {"type": "string", "description": "Файл .docx"},
                  "image_path": {"type": "string"},
                  "width_cm": {"type": "number"}},
              "required": ["path", "image_path"]},
             docx_add_image),
        Tool("xlsx_create",
             "Создать новую книгу Excel из JSON-объекта "
             '{"Имя листа": [[значения строки], ...], ...}. Перезаписывает '
             "файл, если он уже существует.",
             {"type": "object",
              "properties": {
                  "path": {"type": "string"},
                  "sheets_json": {"type": "string"}},
              "required": ["path", "sheets_json"]},
             xlsx_create),
        Tool("xlsx_write_rows",
             "Записать строки (JSON-список списков) в лист книги Excel: "
             "создаёт книгу/лист при отсутствии, по умолчанию дописывает "
             "после последней заполненной строки (start_row=0).",
             {"type": "object",
              "properties": {
                  "path": {"type": "string"},
                  "sheet": {"type": "string"},
                  "rows_json": {"type": "string", "description": '[["A","B"],[1,2]]'},
                  "start_row": {"type": "integer",
                                "description": "0 = дописать после конца листа"},
                  "header_bold": {"type": "boolean",
                                  "description": "Выделить первую записанную строку жирным"}},
              "required": ["path", "sheet", "rows_json"]},
             xlsx_write_rows),
        Tool("xlsx_add_chart",
             "Добавить диаграмму (bar/line/pie) на лист Excel по уже "
             "записанным данным. data_range/categories_range — диапазоны "
             "вида 'B1:B5'.",
             {"type": "object",
              "properties": {
                  "path": {"type": "string"},
                  "sheet": {"type": "string"},
                  "chart_type": {"type": "string", "description": "bar|line|pie"},
                  "data_range": {"type": "string"},
                  "categories_range": {"type": "string"},
                  "title": {"type": "string"},
                  "anchor": {"type": "string", "description": "Ячейка привязки, напр. 'E2'"}},
              "required": ["path", "sheet", "chart_type", "data_range"]},
             xlsx_add_chart),
        Tool("pptx_create",
             "Создать новую презентацию PowerPoint с титульным слайдом. "
             "Перезаписывает файл, если он уже существует.",
             {"type": "object",
              "properties": {
                  "path": {"type": "string"},
                  "title": {"type": "string"},
                  "subtitle": {"type": "string"}},
              "required": ["path"]},
             pptx_create),
        Tool("pptx_add_slide",
             "Добавить слайд с заголовком и телом (markdown: абзацы/списки) "
             "в презентацию; если файла нет — создаёт новый. layout: "
             "title_content|title_only|section_header|blank.",
             {"type": "object",
              "properties": {
                  "path": {"type": "string"},
                  "title": {"type": "string"},
                  "markdown_body": {"type": "string"},
                  "notes": {"type": "string", "description": "Заметки докладчика"},
                  "layout": {"type": "string"}},
              "required": ["path"]},
             pptx_add_slide),
        Tool("pptx_add_table_slide",
             "Добавить слайд с таблицей (JSON header/rows) в презентацию.",
             {"type": "object",
              "properties": {
                  "path": {"type": "string"},
                  "title": {"type": "string"},
                  "header_json": {"type": "string"},
                  "rows_json": {"type": "string"}},
              "required": ["path", "title", "header_json"]},
             pptx_add_table_slide),
        Tool("pptx_add_image_slide",
             "Добавить слайд с изображением (из рабочей папки) в презентацию.",
             {"type": "object",
              "properties": {
                  "path": {"type": "string"},
                  "image_path": {"type": "string"},
                  "title": {"type": "string"},
                  "notes": {"type": "string"}},
              "required": ["path", "image_path"]},
             pptx_add_image_slide),
    ]
