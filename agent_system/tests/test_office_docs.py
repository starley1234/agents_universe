"""Тесты навыка office: создание Word/Excel/PowerPoint.

Требуют python-docx, openpyxl, python-pptx. Если их нет — тесты
пропускаются с понятным сообщением, а не падают с ImportError.

Философия та же, что у test_pdf.py/test_docparse.py: тест обязан уметь
падать. Рядом с позитивными проверками — негативные: невалидный JSON,
неизвестный стиль/раскладка, отсутствующий файл изображения, попытка
добавить диаграмму в несуществующую книгу.
"""
from __future__ import annotations

import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from agent.tools.base import ToolError, Workspace         # noqa: E402

PASS, FAIL = 0, 0

try:
    import docx  # type: ignore
    import openpyxl  # type: ignore
    import pptx  # type: ignore
    HAVE_DEPS = True
except ImportError:
    HAVE_DEPS = False


def check(name: str, cond: bool, detail: str = "") -> None:
    global PASS, FAIL
    if cond:
        PASS += 1
        print(f"  ok   {name}")
    else:
        FAIL += 1
        print(f"  FAIL {name}" + (f" — {detail}" if detail else ""))


def section(title: str) -> None:
    print(f"\n{title}\n" + "─" * len(title))


def _import_office():
    from agent.tools import office_docs
    return office_docs


def _write_test_png(path: Path) -> None:
    """Валидный крошечный PNG для тестов — генерируем через PIL, чтобы не
    хардкодить сырые байты формата (легко ошибиться и получить битый файл,
    который сам ничего не тестирует, кроме обработчика ошибок docx)."""
    from PIL import Image
    Image.new("RGB", (4, 4), (200, 0, 0)).save(path, format="PNG")


MD_SAMPLE = """# Заголовок

Обычный **жирный** и *курсив* текст, `код` тоже.

## Подраздел

- пункт один
- пункт два
  - вложенный пункт
1. первый
2. второй

| Поз | Название |
|---|---|
| 1 | Корпус |
| 2 | Крышка |
"""


# ================================================================ tests
def test_markdown_parser() -> None:
    section("Разбор упрощённого markdown в блоки")
    office = _import_office()
    blocks = office.parse_markdown_blocks(MD_SAMPLE)
    kinds = [b["kind"] for b in blocks]
    check("заголовок 1 уровня распознан", kinds[0] == "heading")
    check("абзац с inline-разметкой распознан", "paragraph" in kinds)
    check("заголовок 2 уровня распознан",
          any(b["kind"] == "heading" and b["level"] == 2 for b in blocks))
    check("маркированный список распознан", "bullet" in kinds)
    check("нумерованный список распознан", "numbered" in kinds)
    check("таблица распознана", "table" in kinds)

    table = next(b for b in blocks if b["kind"] == "table")
    check("заголовок таблицы верный", table["header"] == ["Поз", "Название"])
    check("строки таблицы верные", table["rows"] == [["1", "Корпус"], ["2", "Крышка"]])

    bullet = next(b for b in blocks if b["kind"] == "bullet")
    check("вложенность списка учтена",
          any(lvl == 1 for lvl, _ in bullet["items"]), str(bullet["items"]))


def test_docx_create_and_append() -> None:
    section("Word: создание, дополнение, таблицы, изображения")
    office = _import_office()
    with tempfile.TemporaryDirectory() as td:
        ws = Workspace(Path(td) / "ws")
        tools = {t.name: t for t in office.build(ws)}

        out = tools["docx_create"].fn(path="doc.docx", title="Титул",
                                      markdown=MD_SAMPLE)
        check("документ создан", "Создан" in out, out)

        out2 = tools["docx_append"].fn(path="doc.docx", markdown="Ещё текст.")
        check("документ дополнен", "Дополнен" in out2, out2)

        out3 = tools["docx_add_table"].fn(
            path="doc.docx", header_json='["A","B"]', rows_json='[["1","2"]]')
        check("таблица добавлена отдельным инструментом", "Таблица добавлена" in out3)

        d = docx.Document(str(ws.root / "doc.docx"))
        texts = [p.text for p in d.paragraphs]
        check("заголовок титульного уровня на месте", "Титул" in texts)
        check("заголовок 1 уровня перенесён", "Заголовок" in texts)
        check("вложенный список получил свой стиль",
              any(p.style.name == "List Bullet 2" for p in d.paragraphs))
        check("дополнение действительно в конце", texts[-1] == "Ещё текст.")
        check("обе таблицы на месте (из markdown и docx_add_table)",
              len(d.tables) == 2)
        check("данные первой таблицы верны",
              d.tables[0].rows[1].cells[1].text == "Корпус")


def test_docx_image() -> None:
    section("Word: вставка изображения")
    office = _import_office()
    with tempfile.TemporaryDirectory() as td:
        ws = Workspace(Path(td) / "ws")
        tools = {t.name: t for t in office.build(ws)}
        _write_test_png(ws.root / "pic.png")
        out = tools["docx_add_image"].fn(path="doc.docx", image_path="pic.png")
        check("изображение добавлено", "добавлено" in out, out)
        d = docx.Document(str(ws.root / "doc.docx"))
        check("в документе есть картинка", len(d.inline_shapes) == 1)



def test_xlsx_create_write_chart() -> None:
    section("Excel: листы, дозапись строк, диаграмма")
    office = _import_office()
    with tempfile.TemporaryDirectory() as td:
        ws = Workspace(Path(td) / "ws")
        tools = {t.name: t for t in office.build(ws)}

        out = tools["xlsx_create"].fn(
            path="book.xlsx",
            sheets_json='{"Data": [["A","B"],[1,2]], "Notes": [["текст"]]}')
        check("книга создана с 2 листами", "2 лист" in out, out)

        out2 = tools["xlsx_write_rows"].fn(
            path="book.xlsx", sheet="Data", rows_json="[[3,4]]")
        check("строка дозаписана после конца", "строки 3" in out2, out2)

        wb = openpyxl.load_workbook(str(ws.root / "book.xlsx"))
        rows = list(wb["Data"].iter_rows(values_only=True))
        check("данные на месте", rows == [("A", "B"), (1, 2), (3, 4)], str(rows))

        out3 = tools["xlsx_add_chart"].fn(
            path="book.xlsx", sheet="Data", chart_type="bar",
            data_range="B1:B3", categories_range="A2:A3", title="Тест")
        check("диаграмма добавлена", "Диаграмма bar" in out3, out3)


def test_pptx_create_slides() -> None:
    section("PowerPoint: слайды, таблица, заметки докладчика")
    office = _import_office()
    with tempfile.TemporaryDirectory() as td:
        ws = Workspace(Path(td) / "ws")
        tools = {t.name: t for t in office.build(ws)}

        out = tools["pptx_create"].fn(path="p.pptx", title="Заголовок",
                                      subtitle="Подзаголовок")
        check("презентация создана", "1 слайд" in out, out)

        out2 = tools["pptx_add_slide"].fn(
            path="p.pptx", title="Слайд 2",
            markdown_body="- пункт один\n- пункт два", notes="заметка докладчика")
        check("слайд добавлен", "2 слайд" in out2, out2)

        out3 = tools["pptx_add_table_slide"].fn(
            path="p.pptx", title="Таблица", header_json='["X","Y"]',
            rows_json='[["1","2"],["3","4"]]')
        check("слайд с таблицей добавлен", "2 строк" in out3, out3)

        prs = pptx.Presentation(str(ws.root / "p.pptx"))
        check("всего 3 слайда", len(prs.slides) == 3, str(len(prs.slides)))
        check("заголовок титульного слайда верный",
              prs.slides[0].shapes.title.text == "Заголовок")
        check("заметки докладчика сохранены",
              prs.slides[1].notes_slide.notes_text_frame.text == "заметка докладчика")

        table_shape = next(s for s in prs.slides[2].shapes if s.has_table)
        check("данные таблицы на слайде верны",
              table_shape.table.cell(1, 1).text == "2")


def test_negative_cases() -> None:
    section("Негативные проверки")
    office = _import_office()
    with tempfile.TemporaryDirectory() as td:
        ws = Workspace(Path(td) / "ws")
        tools = {t.name: t for t in office.build(ws)}

        try:
            tools["xlsx_create"].fn(path="x.xlsx", sheets_json="не json")
            check("отказ на невалидный JSON листов", False)
        except ToolError:
            check("отказ на невалидный JSON листов", True)

        try:
            tools["xlsx_add_chart"].fn(path="nope.xlsx", sheet="s",
                                       chart_type="bar", data_range="A1:A2")
            check("отказ на отсутствующую книгу", False)
        except ToolError:
            check("отказ на отсутствующую книгу", True)

        tools["xlsx_create"].fn(path="b.xlsx", sheets_json='{"S": [["a"]]}')
        try:
            tools["xlsx_add_chart"].fn(path="b.xlsx", sheet="Нет такого",
                                       chart_type="bar", data_range="A1:A2")
            check("отказ на отсутствующий лист", False)
        except ToolError:
            check("отказ на отсутствующий лист", True)

        try:
            tools["xlsx_add_chart"].fn(path="b.xlsx", sheet="S",
                                       chart_type="выдуманный", data_range="A1:A2")
            check("отказ на неизвестный тип диаграммы", False)
        except ToolError:
            check("отказ на неизвестный тип диаграммы", True)

        try:
            tools["docx_add_table"].fn(path="d.docx", header_json='["A"]',
                                       rows_json="[]", style="ФейкСтиль")
            check("отказ на неизвестный стиль таблицы", False)
        except ToolError:
            check("отказ на неизвестный стиль таблицы", True)

        try:
            tools["pptx_add_slide"].fn(path="p.pptx", layout="выдуманное")
            check("отказ на неизвестную раскладку", False)
        except ToolError:
            check("отказ на неизвестную раскладку", True)

        try:
            tools["docx_add_image"].fn(path="d.docx", image_path="нет_такого.png")
            check("отказ на отсутствующее изображение", False)
        except ToolError:
            check("отказ на отсутствующее изображение", True)

        try:
            tools["xlsx_write_rows"].fn(path="b.xlsx", sheet="S", rows_json="{}")
            check("отказ, если rows_json не список списков", False)
        except ToolError:
            check("отказ, если rows_json не список списков", True)


def test_build_agent_with_office_skill() -> None:
    section("Сборка агента с навыком office")
    from agent.build import build_agent
    from agent.config import Config
    with tempfile.TemporaryDirectory() as td:
        cfg = Config(provider="ollama", model="m", workspace=td,
                    skills=["files", "office"])
        agent = build_agent(cfg)
        names = agent.tools.names()
        for t in ("docx_create", "docx_append", "docx_add_table", "docx_add_image",
                  "xlsx_create", "xlsx_write_rows", "xlsx_add_chart",
                  "pptx_create", "pptx_add_slide", "pptx_add_table_slide",
                  "pptx_add_image_slide"):
            check(f"инструмент {t} зарегистрирован", t in names)


def main() -> int:
    if not HAVE_DEPS:
        print("python-docx/openpyxl/python-pptx не установлены — тесты office "
              "пропущены (pip install python-docx openpyxl python-pptx)")
        return 0
    test_markdown_parser()
    test_docx_create_and_append()
    test_docx_image()
    test_xlsx_create_write_chart()
    test_pptx_create_slides()
    test_negative_cases()
    test_build_agent_with_office_skill()

    print(f"\n{'─' * 40}\nитого: {PASS} ok, {FAIL} fail")
    return 1 if FAIL else 0


if __name__ == "__main__":
    raise SystemExit(main())
