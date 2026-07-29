"""Создание документов: Word, Excel, PowerPoint, PDF.

Один вход на все форматы — Markdown. Причина простая: заставлять модель
помнить четыре разных набора вызовов значит получать четыре набора
ошибок. Агент пишет привычную разметку, конвертер разворачивает её в
нужный формат.

    # Заголовок 1        -> заголовок / новый слайд
    ## Заголовок 2       -> подзаголовок
    обычный текст        -> абзац
    - пункт              -> список
    | а | б |            -> таблица
    ---                  -> разрыв слайда/страницы

Библиотеки используются, если есть (python-docx, openpyxl, python-pptx):
качество выше. Нет — работает запасной путь на стандартной библиотеке,
потому что docx/xlsx/pptx это ZIP с XML внутри. Обязательных
зависимостей по-прежнему ноль.

PDF собирается сам, без библиотек, и это единственное место со сложным
кодом: чтобы кириллица не превратилась в вопросительные знаки, шрифт
TrueType встраивается в файл вместе с картой ToUnicode. Без неё текст
нельзя ни найти, ни скопировать из готового PDF.
"""
from __future__ import annotations

import html
import re
import struct
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Callable

from ..store import Store
from ..tools.base import Tool, ToolError, Workspace

#: Шрифты с кириллицей, в порядке предпочтения. Первый найденный
#: встраивается в PDF.
FONT_PATHS = [
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    "/usr/share/fonts/TTF/DejaVuSans.ttf",
    "/Library/Fonts/Arial.ttf",
]


# ═══════════════════════ разбор Markdown ═══════════════════════════
@dataclass
class Block:
    """Кусок документа. Пять видов покрывают деловые бумаги целиком."""
    kind: str                    # heading | text | list | table | break
    level: int = 0               # для heading
    text: str = ""
    items: list[str] = field(default_factory=list)
    rows: list[list[str]] = field(default_factory=list)


_TABLE_SEP = re.compile(r"^\s*\|?[\s:|-]+\|[\s:|-]*$")


def parse_markdown(src: str) -> list[Block]:
    """Markdown -> плоский список блоков.

    Намеренно понимаем немного: заголовки, абзацы, списки, таблицы,
    разрыв. Полный Markdown здесь не нужен, а его поддержка утянула бы
    за собой разбор ссылок, кода и вложенности.
    """
    blocks: list[Block] = []
    lines = src.replace("\r\n", "\n").split("\n")
    i, para, items = 0, [], []

    def flush() -> None:
        nonlocal para, items
        if para:
            blocks.append(Block("text", text=" ".join(para).strip()))
            para = []
        if items:
            blocks.append(Block("list", items=items))
            items = []

    while i < len(lines):
        ln = lines[i]
        s = ln.strip()

        if not s:
            flush()
            i += 1
            continue

        # разрыв: --- или ***
        if re.fullmatch(r"(-{3,}|\*{3,})", s):
            flush()
            blocks.append(Block("break"))
            i += 1
            continue

        # заголовок
        m = re.match(r"^(#{1,6})\s+(.*)$", s)
        if m:
            flush()
            blocks.append(Block("heading", level=len(m.group(1)),
                                text=_plain(m.group(2))))
            i += 1
            continue

        # таблица: строка с | и следующая — разделитель
        if "|" in s and i + 1 < len(lines) and _TABLE_SEP.match(lines[i + 1]):
            flush()
            rows = []
            head = _split_row(s)
            rows.append(head)
            i += 2
            while i < len(lines) and "|" in lines[i]:
                rows.append(_split_row(lines[i]))
                i += 1
            width = max(len(r) for r in rows)
            rows = [r + [""] * (width - len(r)) for r in rows]
            blocks.append(Block("table", rows=rows))
            continue

        # список
        m = re.match(r"^[-*+•]\s+(.*)$", s) or re.match(r"^\d+[.)]\s+(.*)$", s)
        if m:
            if para:
                flush()
            items.append(_plain(m.group(1)))
            i += 1
            continue

        if items:
            flush()
        para.append(_plain(s))
        i += 1

    flush()
    return blocks


def _split_row(line: str) -> list[str]:
    cells = line.strip().strip("|").split("|")
    return [_plain(c.strip()) for c in cells]


def _plain(t: str) -> str:
    """Снять разметку выделения: в деловом документе она лишний шум."""
    t = re.sub(r"\*\*(.+?)\*\*", r"\1", t)
    t = re.sub(r"__(.+?)__", r"\1", t)
    t = re.sub(r"(?<!\*)\*([^*]+)\*(?!\*)", r"\1", t)
    t = re.sub(r"`([^`]+)`", r"\1", t)
    t = re.sub(r"\[([^\]]+)\]\([^)]*\)", r"\1", t)   # ссылку -> текст
    return t.strip()


# ═════════════════════════════ DOCX ════════════════════════════════
W = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"


def _docx_lib(path: Path, blocks: list[Block], title: str) -> bool:
    try:
        import docx                                      # type: ignore
    except ImportError:
        return False
    d = docx.Document()
    if title:
        d.add_heading(title, 0)
    for b in blocks:
        if b.kind == "heading":
            d.add_heading(b.text, min(b.level, 9))
        elif b.kind == "text":
            d.add_paragraph(b.text)
        elif b.kind == "list":
            for it in b.items:
                d.add_paragraph(it, style="List Bullet")
        elif b.kind == "table" and b.rows:
            t = d.add_table(rows=len(b.rows), cols=len(b.rows[0]))
            t.style = "Table Grid"
            for ri, row in enumerate(b.rows):
                for ci, cell in enumerate(row):
                    t.cell(ri, ci).text = cell
        elif b.kind == "break":
            d.add_page_break()
    d.save(str(path))
    return True


def _docx_plain(path: Path, blocks: list[Block], title: str) -> None:
    """Запасной путь: docx это ZIP с word/document.xml."""
    body: list[str] = []

    def para(text: str, style: str = "") -> str:
        pr = f'<w:pPr><w:pStyle w:val="{style}"/></w:pPr>' if style else ""
        return (f"<w:p>{pr}<w:r><w:t xml:space='preserve'>"
                f"{html.escape(text)}</w:t></w:r></w:p>")

    if title:
        body.append(para(title, "Title"))
    for b in blocks:
        if b.kind == "heading":
            body.append(para(b.text, f"Heading{min(b.level, 6)}"))
        elif b.kind == "text":
            body.append(para(b.text))
        elif b.kind == "list":
            for it in b.items:
                body.append(para("• " + it))
        elif b.kind == "table" and b.rows:
            cells = []
            for row in b.rows:
                tcs = "".join(
                    "<w:tc><w:tcPr><w:tcBorders>"
                    + "".join(f'<w:{e} w:val="single" w:sz="4"/>'
                              for e in ("top", "left", "bottom", "right"))
                    + "</w:tcBorders></w:tcPr>"
                    + para(c) + "</w:tc>" for c in row)
                cells.append(f"<w:tr>{tcs}</w:tr>")
            body.append("<w:tbl>" + "".join(cells) + "</w:tbl>")
        elif b.kind == "break":
            body.append('<w:p><w:r><w:br w:type="page"/></w:r></w:p>')

    doc = (f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
           f'<w:document xmlns:w="{W}"><w:body>{"".join(body)}</w:body>'
           f'</w:document>')
    rels = ('<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/'
            'package/2006/relationships"><Relationship Id="rId1" '
            'Type="http://schemas.openxmlformats.org/officeDocument/2006/'
            'relationships/officeDocument" Target="word/document.xml"/>'
            '</Relationships>')
    ct = ('<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
          '<Types xmlns="http://schemas.openxmlformats.org/package/2006/'
          'content-types"><Default Extension="rels" ContentType='
          '"application/vnd.openxmlformats-package.relationships+xml"/>'
          '<Default Extension="xml" ContentType="application/xml"/>'
          '<Override PartName="/word/document.xml" ContentType='
          '"application/vnd.openxmlformats-officedocument.'
          'wordprocessingml.document.main+xml"/></Types>')
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", ct)
        z.writestr("_rels/.rels", rels)
        z.writestr("word/document.xml", doc)


# ═════════════════════════════ XLSX ════════════════════════════════
def _sheets_from_blocks(blocks: list[Block]) -> list[tuple[str, list[list[str]]]]:
    """Таблицы -> листы. Имя листа берём из ближайшего заголовка выше."""
    out: list[tuple[str, list[list[str]]]] = []
    name = ""
    for b in blocks:
        if b.kind == "heading":
            name = b.text
        elif b.kind == "table" and b.rows:
            title = _sheet_name(name or f"Лист{len(out) + 1}",
                                [n for n, _ in out])
            out.append((title, b.rows))
            name = ""
    return out


def _sheet_name(raw: str, taken: list[str]) -> str:
    """Excel не принимает : \\ / ? * [ ] и длину больше 31."""
    n = re.sub(r"[:\\/?*\[\]]", "-", raw).strip() or "Лист"
    n = n[:31]
    base, k = n, 2
    while n in taken:
        suf = f"_{k}"
        n = base[:31 - len(suf)] + suf
        k += 1
    return n


def _xlsx_lib(path: Path, sheets: list[tuple[str, list[list[str]]]]) -> bool:
    try:
        import openpyxl                                  # type: ignore
    except ImportError:
        return False
    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    for name, rows in sheets:
        ws = wb.create_sheet(title=name)
        for row in rows:
            ws.append([_as_number(c) for c in row])
        # ширина по содержимому: иначе всё слипается в узкие столбцы
        for ci in range(1, len(rows[0]) + 1):
            width = max((len(str(r[ci - 1])) for r in rows if ci <= len(r)),
                        default=8)
            ws.column_dimensions[
                openpyxl.utils.get_column_letter(ci)].width = min(width + 2, 60)
        if len(rows) > 1:
            ws.freeze_panes = "A2"
    wb.save(str(path))
    return True


def _as_number(cell: str) -> Any:
    """Число должно попасть в ячейку числом, иначе не посчитается сумма."""
    s = str(cell).strip().replace("\u00a0", "").replace(" ", "")
    if not s:
        return ""
    t = s.replace(",", ".") if s.count(",") == 1 and "." not in s else s
    try:
        if re.fullmatch(r"[-+]?\d+", t):
            return int(t)
        if re.fullmatch(r"[-+]?\d*\.\d+([eE][-+]?\d+)?", t):
            return float(t)
    except ValueError:
        pass
    return cell


def _xlsx_plain(path: Path, sheets: list[tuple[str, list[list[str]]]]) -> None:
    """Запасной путь. Строки пишем как inlineStr: без общей таблицы строк
    файл проще, а Excel такое читает."""
    NS = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
    R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    def col(i: int) -> str:
        s = ""
        while i >= 0:
            s = chr(ord("A") + i % 26) + s
            i = i // 26 - 1
        return s

    parts, entries, rels = [], [], []
    for k, (name, rows) in enumerate(sheets, 1):
        xml_rows = []
        for ri, row in enumerate(rows, 1):
            cs = []
            for ci, cell in enumerate(row):
                v = _as_number(cell)
                ref = f"{col(ci)}{ri}"
                if isinstance(v, (int, float)):
                    cs.append(f'<c r="{ref}"><v>{v}</v></c>')
                elif str(cell).strip():
                    cs.append(f'<c r="{ref}" t="inlineStr"><is><t '
                              f'xml:space="preserve">{html.escape(str(cell))}'
                              f'</t></is></c>')
            xml_rows.append(f'<row r="{ri}">{"".join(cs)}</row>')
        parts.append((f"xl/worksheets/sheet{k}.xml",
                      f'<?xml version="1.0" encoding="UTF-8"?>'
                      f'<worksheet xmlns="{NS}"><sheetData>'
                      f'{"".join(xml_rows)}</sheetData></worksheet>'))
        entries.append(f'<sheet name="{html.escape(name)}" sheetId="{k}" '
                       f'r:id="rId{k}"/>')
        rels.append(f'<Relationship Id="rId{k}" Type="{R}/worksheet" '
                    f'Target="worksheets/sheet{k}.xml"/>')

    wb = (f'<?xml version="1.0" encoding="UTF-8"?><workbook xmlns="{NS}" '
          f'xmlns:r="{R}"><sheets>{"".join(entries)}</sheets></workbook>')
    ct_over = "".join(
        f'<Override PartName="/xl/worksheets/sheet{k}.xml" ContentType='
        f'"application/vnd.openxmlformats-officedocument.spreadsheetml.'
        f'worksheet+xml"/>' for k in range(1, len(sheets) + 1))
    ct = ('<?xml version="1.0" encoding="UTF-8"?><Types xmlns='
          '"http://schemas.openxmlformats.org/package/2006/content-types">'
          '<Default Extension="rels" ContentType="application/vnd.'
          'openxmlformats-package.relationships+xml"/><Default '
          'Extension="xml" ContentType="application/xml"/><Override '
          'PartName="/xl/workbook.xml" ContentType="application/vnd.'
          'openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>'
          + ct_over + "</Types>")
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", ct)
        z.writestr("_rels/.rels",
                   '<?xml version="1.0" encoding="UTF-8"?><Relationships '
                   'xmlns="http://schemas.openxmlformats.org/package/2006/'
                   f'relationships"><Relationship Id="rIdWb" Type="{R}/'
                   'officeDocument" Target="xl/workbook.xml"/></Relationships>')
        z.writestr("xl/workbook.xml", wb)
        z.writestr("xl/_rels/workbook.xml.rels",
                   '<?xml version="1.0" encoding="UTF-8"?><Relationships '
                   'xmlns="http://schemas.openxmlformats.org/package/2006/'
                   f'relationships">{"".join(rels)}</Relationships>')
        for name, data in parts:
            z.writestr(name, data)


# ═════════════════════════════ PPTX ════════════════════════════════
def _slides_from_blocks(blocks: list[Block]) -> list[dict[str, Any]]:
    """Блоки -> слайды. Новый слайд начинает заголовок или разрыв."""
    slides: list[dict[str, Any]] = []
    cur: dict[str, Any] = {"title": "", "bullets": [], "rows": []}

    def push() -> None:
        if cur["title"] or cur["bullets"] or cur["rows"]:
            slides.append(dict(cur))
        cur["title"], cur["bullets"], cur["rows"] = "", [], []

    for b in blocks:
        if b.kind == "heading" and b.level <= 2:
            push()
            cur["title"] = b.text
        elif b.kind == "heading":
            cur["bullets"].append(b.text)
        elif b.kind == "break":
            push()
        elif b.kind == "text":
            cur["bullets"].append(b.text)
        elif b.kind == "list":
            cur["bullets"].extend(b.items)
        elif b.kind == "table":
            cur["rows"] = b.rows
    push()
    return slides


def _pptx_lib(path: Path, slides: list[dict[str, Any]]) -> bool:
    try:
        from pptx import Presentation                    # type: ignore
        from pptx.util import Inches, Pt                 # type: ignore
    except ImportError:
        return False
    prs = Presentation()
    for s in slides:
        if s["rows"]:
            sl = prs.slides.add_slide(prs.slide_layouts[5])
            sl.shapes.title.text = s["title"] or " "
            rows, cols = len(s["rows"]), len(s["rows"][0])
            shape = sl.shapes.add_table(rows, cols, Inches(0.5), Inches(1.8),
                                        Inches(9), Inches(0.4 * rows))
            for ri, row in enumerate(s["rows"]):
                for ci, cell in enumerate(row):
                    tc = shape.table.cell(ri, ci)
                    tc.text = cell
                    for p in tc.text_frame.paragraphs:
                        for r in p.runs:
                            r.font.size = Pt(12)
        else:
            sl = prs.slides.add_slide(prs.slide_layouts[1])
            sl.shapes.title.text = s["title"] or " "
            body = sl.placeholders[1].text_frame
            body.clear()
            for i, b in enumerate(s["bullets"]):
                p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                p.text = b
    prs.save(str(path))
    return True


_P_NS = ("xmlns:a='http://schemas.openxmlformats.org/drawingml/2006/main' "
         "xmlns:p='http://schemas.openxmlformats.org/presentationml/2006/main'")


def _pptx_plain(path: Path, slides: list[dict[str, Any]]) -> None:
    """Запасной путь: минимальная презентация, текст в надписях."""
    def tx(text: str, x: int, y: int, cx: int, cy: int, size: int,
           idx: int) -> str:
        paras = "".join(
            f"<a:p><a:r><a:rPr lang='ru-RU' sz='{size}'/>"
            f"<a:t>{html.escape(ln)}</a:t></a:r></a:p>"
            for ln in (text.split("\n") if text else [" "]))
        return (f"<p:sp><p:nvSpPr><p:cNvPr id='{idx}' name='t{idx}'/>"
                f"<p:cNvSpPr txBox='1'/><p:nvPr/></p:nvSpPr>"
                f"<p:spPr><a:xfrm><a:off x='{x}' y='{y}'/>"
                f"<a:ext cx='{cx}' cy='{cy}'/></a:xfrm>"
                f"<a:prstGeom prst='rect'><a:avLst/></a:prstGeom></p:spPr>"
                f"<p:txBody><a:bodyPr wrap='square'/><a:lstStyle/>"
                f"{paras}</p:txBody></p:sp>")

    EMU = 914400
    slide_xml = []
    for s in slides:
        shapes = tx(s["title"], EMU // 2, EMU // 3, 8 * EMU, EMU, 3200, 2)
        lines = list(s["bullets"])
        for row in s["rows"]:
            lines.append(" | ".join(row))
        if lines:
            body = "\n".join(("• " + ln) if ln in s["bullets"] else ln
                             for ln in lines)
            shapes += tx(body, EMU // 2, EMU * 3 // 2, 8 * EMU, 4 * EMU,
                         1800, 3)
        slide_xml.append(
            f"<?xml version='1.0' encoding='UTF-8'?><p:sld {_P_NS}>"
            f"<p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id='1' name=''/>"
            f"<p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/>"
            f"{shapes}</p:spTree></p:cSld></p:sld>")

    R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    ids = "".join(f"<p:sldId id='{255 + i}' r:id='rId{i + 1}'/>"
                  for i in range(len(slide_xml)))
    pres = (f"<?xml version='1.0' encoding='UTF-8'?><p:presentation {_P_NS} "
            f"xmlns:r='{R}'><p:sldIdLst>{ids}</p:sldIdLst>"
            f"<p:sldSz cx='9144000' cy='6858000'/>"
            f"<p:notesSz cx='6858000' cy='9144000'/></p:presentation>")
    prels = "".join(f"<Relationship Id='rId{i + 1}' Type='{R}/slide' "
                    f"Target='slides/slide{i + 1}.xml'/>"
                    for i in range(len(slide_xml)))
    ct = ('<?xml version="1.0" encoding="UTF-8"?><Types xmlns='
          '"http://schemas.openxmlformats.org/package/2006/content-types">'
          '<Default Extension="rels" ContentType="application/vnd.'
          'openxmlformats-package.relationships+xml"/><Default '
          'Extension="xml" ContentType="application/xml"/><Override '
          'PartName="/ppt/presentation.xml" ContentType="application/vnd.'
          'openxmlformats-officedocument.presentationml.presentation.main'
          '+xml"/>' + "".join(
              f'<Override PartName="/ppt/slides/slide{i + 1}.xml" '
              f'ContentType="application/vnd.openxmlformats-officedocument.'
              f'presentationml.slide+xml"/>' for i in range(len(slide_xml)))
          + "</Types>")
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", ct)
        z.writestr("_rels/.rels",
                   '<?xml version="1.0" encoding="UTF-8"?><Relationships '
                   'xmlns="http://schemas.openxmlformats.org/package/2006/'
                   f'relationships"><Relationship Id="rIdP" Type="{R}/'
                   'officeDocument" Target="ppt/presentation.xml"/>'
                   '</Relationships>')
        z.writestr("ppt/presentation.xml", pres)
        z.writestr("ppt/_rels/presentation.xml.rels",
                   '<?xml version="1.0" encoding="UTF-8"?><Relationships '
                   'xmlns="http://schemas.openxmlformats.org/package/2006/'
                   f'relationships">{prels}</Relationships>')
        for i, x in enumerate(slide_xml, 1):
            z.writestr(f"ppt/slides/slide{i}.xml", x)


# ══════════════════════════════ PDF ════════════════════════════════
class TrueType:
    """Минимальный разбор TTF: карта символов и ширины.

    Нужен ровно для одного: встроить шрифт в PDF так, чтобы кириллица
    отображалась и текст можно было выделить. Без встраивания в PDF
    доступны только латинские базовые шрифты.
    """

    def __init__(self, path: Path) -> None:
        self.data = path.read_bytes()
        self.tables: dict[str, tuple[int, int]] = {}
        num = struct.unpack(">H", self.data[4:6])[0]
        for i in range(num):
            off = 12 + i * 16
            tag = self.data[off:off + 4].decode("latin-1")
            start, length = struct.unpack(">II", self.data[off + 8:off + 16])
            self.tables[tag] = (start, length)
        head = self.tables["head"][0]
        self.units = struct.unpack(">H", self.data[head + 18:head + 20])[0]
        self.index_to_loc = struct.unpack(
            ">h", self.data[head + 50:head + 52])[0]
        hhea = self.tables["hhea"][0]
        self.num_hmetrics = struct.unpack(
            ">H", self.data[hhea + 34:hhea + 36])[0]
        self.cmap = self._read_cmap()
        self._widths: dict[int, int] = {}

    def _read_cmap(self) -> dict[int, int]:
        start = self.tables["cmap"][0]
        n = struct.unpack(">H", self.data[start + 2:start + 4])[0]
        best = 0
        for i in range(n):
            off = start + 4 + i * 8
            pid, eid, sub = struct.unpack(">HHI", self.data[off:off + 8])
            if (pid, eid) in ((3, 1), (3, 10), (0, 3), (0, 4)):
                best = start + sub
                break
        if not best:
            raise ToolError("в шрифте нет пригодной таблицы символов")
        fmt = struct.unpack(">H", self.data[best:best + 2])[0]
        if fmt != 4:
            raise ToolError(f"формат cmap {fmt} не поддержан")
        segx2 = struct.unpack(">H", self.data[best + 6:best + 8])[0]
        seg = segx2 // 2
        base = best + 14

        def arr(o: int) -> list[int]:
            return list(struct.unpack(f">{seg}H",
                                      self.data[o:o + segx2]))
        ends = arr(base)
        starts = arr(base + segx2 + 2)
        deltas = arr(base + segx2 * 2 + 2)
        range_off_pos = base + segx2 * 3 + 2
        ranges = arr(range_off_pos)

        out: dict[int, int] = {}
        for i in range(seg):
            for c in range(starts[i], min(ends[i], 0xFFFF) + 1):
                if ranges[i] == 0:
                    g = (c + deltas[i]) & 0xFFFF
                else:
                    p = range_off_pos + i * 2 + ranges[i] + (c - starts[i]) * 2
                    if p + 2 > len(self.data):
                        continue
                    g = struct.unpack(">H", self.data[p:p + 2])[0]
                    if g:
                        g = (g + deltas[i]) & 0xFFFF
                if g:
                    out[c] = g
        return out

    def width(self, gid: int) -> int:
        """Ширина глифа в 1/1000 em — единицах PDF."""
        if gid in self._widths:
            return self._widths[gid]
        hmtx = self.tables["hmtx"][0]
        i = min(gid, self.num_hmetrics - 1)
        raw = struct.unpack(">H", self.data[hmtx + i * 4:hmtx + i * 4 + 2])[0]
        w = round(raw * 1000 / self.units)
        self._widths[gid] = w
        return w

    def gid(self, ch: str) -> int:
        return self.cmap.get(ord(ch), 0)


def _pdf_escape(b: bytes) -> bytes:
    return b.replace(b"\\", b"\\\\").replace(b"(", b"\\(").replace(b")", b"\\)")


def _make_pdf(path: Path, blocks: list[Block], title: str,
              font_path: Path) -> None:
    """Собрать PDF со встроенным шрифтом.

    Кодировка Identity-H: в поток пишутся номера глифов, а карта
    ToUnicode возвращает их обратно в символы — без неё текст в готовом
    файле нельзя ни найти поиском, ни скопировать.
    """
    ttf = TrueType(font_path)
    W, H = 595, 842                       # A4 в пунктах
    LEFT, RIGHT, TOP, BOT = 57, 538, 785, 57
    used: set[int] = set()

    def wrap(text: str, size: float, width: float) -> list[str]:
        words, lines, cur = text.split(), [], ""
        for w in words:
            trial = (cur + " " + w).strip()
            if _text_width(ttf, trial) * size / 1000 <= width or not cur:
                cur = trial
            else:
                lines.append(cur)
                cur = w
        if cur:
            lines.append(cur)
        return lines or [""]

    pages: list[list[str]] = []
    ops: list[str] = []
    y = TOP

    def show(text: str, size: float, x: float, dy: float) -> None:
        nonlocal y, ops
        for line in wrap(text, size, RIGHT - x):
            if y - dy < BOT:
                pages.append(ops)
                ops = []
                y = TOP
            y -= dy
            gids = []
            for ch in line:
                g = ttf.gid(ch)
                used.add(g)
                gids.append(f"{g:04X}")
            ops.append(f"BT /F1 {size} Tf {x} {y} Td <{''.join(gids)}> Tj ET")

    if title:
        show(title, 20, LEFT, 30)
        y -= 8
    for b in blocks:
        if b.kind == "heading":
            size = max(11, 18 - b.level * 2)
            y -= 6
            show(b.text, size, LEFT, size + 6)
        elif b.kind == "text":
            show(b.text, 11, LEFT, 15)
        elif b.kind == "list":
            for it in b.items:
                show("• " + it, 11, LEFT + 12, 15)
        elif b.kind == "table":
            for ri, row in enumerate(b.rows):
                show(" | ".join(row), 10, LEFT, 14)
                if ri == 0:
                    y -= 2
        elif b.kind == "break":
            pages.append(ops)
            ops = []
            y = TOP
        y -= 4
    pages.append(ops)
    pages = [p for p in pages if p] or [[]]

    # ---- объекты PDF
    font_data = font_path.read_bytes()
    gids = sorted(used | {0})
    w_parts = " ".join(f"{g} [{ttf.width(g)}]" for g in gids)
    tou = _to_unicode(ttf, gids)

    objs: dict[int, bytes] = {}
    n_pages = len(pages)
    first_page = 5                       # 1 catalog, 2 pages, 3 font, 4 desc
    content_start = first_page + n_pages

    kids = " ".join(f"{first_page + i} 0 R" for i in range(n_pages))
    objs[1] = b"<< /Type /Catalog /Pages 2 0 R >>"
    objs[2] = (f"<< /Type /Pages /Count {n_pages} /Kids [{kids}] >>"
               ).encode()
    objs[3] = (f"<< /Type /Font /Subtype /Type0 /BaseFont /Embedded "
               f"/Encoding /Identity-H /DescendantFonts [{content_start + n_pages} 0 R] "
               f"/ToUnicode {content_start + n_pages + 3} 0 R >>").encode()
    objs[4] = (f"<< /Type /FontDescriptor /FontName /Embedded /Flags 4 "
               f"/FontBBox [-1000 -300 2000 1000] /ItalicAngle 0 /Ascent 900 "
               f"/Descent -200 /CapHeight 700 /StemV 80 "
               f"/FontFile2 {content_start + n_pages + 2} 0 R >>").encode()
    for i in range(n_pages):
        objs[first_page + i] = (
            f"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 {W} {H}] "
            f"/Resources << /Font << /F1 3 0 R >> >> "
            f"/Contents {content_start + i} 0 R >>").encode()
    for i, ops_page in enumerate(pages):
        stream = "\n".join(ops_page).encode("latin-1", "replace")
        objs[content_start + i] = (
            f"<< /Length {len(stream)} >>\nstream\n".encode()
            + stream + b"\nendstream")
    cid = content_start + n_pages
    objs[cid] = (f"<< /Type /Font /Subtype /CIDFontType2 /BaseFont /Embedded "
                 f"/CIDSystemInfo << /Registry (Adobe) /Ordering (Identity) "
                 f"/Supplement 0 >> /FontDescriptor 4 0 R /DW 500 "
                 f"/W [{w_parts}] /CIDToGIDMap /Identity >>").encode()
    objs[cid + 1] = b"<< >>"             # запас, чтобы номера не путались
    objs[cid + 2] = (f"<< /Length {len(font_data)} /Length1 {len(font_data)} >>"
                     f"\nstream\n").encode() + font_data + b"\nendstream"
    objs[cid + 3] = (f"<< /Length {len(tou)} >>\nstream\n".encode()
                     + tou + b"\nendstream")

    out = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets: dict[int, int] = {}
    for num in sorted(objs):
        offsets[num] = len(out)
        out += f"{num} 0 obj\n".encode() + objs[num] + b"\nendobj\n"
    xref_at = len(out)
    top = max(objs) + 1
    out += f"xref\n0 {top}\n".encode()
    out += b"0000000000 65535 f \n"
    for num in range(1, top):
        off = offsets.get(num, 0)
        out += f"{off:010d} 00000 n \n".encode()
    out += (f"trailer\n<< /Size {top} /Root 1 0 R >>\nstartxref\n"
            f"{xref_at}\n%%EOF\n").encode()
    path.write_bytes(bytes(out))


def _text_width(ttf: TrueType, text: str) -> float:
    return sum(ttf.width(ttf.gid(c)) for c in text)


def _to_unicode(ttf: TrueType, gids: list[int]) -> bytes:
    """Карта «глиф -> символ». Без неё текст в PDF не ищется и не копируется."""
    back: dict[int, int] = {}
    for code, g in ttf.cmap.items():
        if g in gids and g not in back:
            back[g] = code
    pairs = [f"<{g:04X}> <{back[g]:04X}>" for g in gids if g in back]
    chunks = []
    for i in range(0, len(pairs), 100):
        part = pairs[i:i + 100]
        chunks.append(f"{len(part)} beginbfchar\n" + "\n".join(part)
                      + "\nendbfchar")
    body = "\n".join(chunks)
    return (
        "/CIDInit /ProcSet findresource begin\n12 dict begin\nbegincmap\n"
        "/CIDSystemInfo << /Registry (Adobe) /Ordering (UCS) /Supplement 0 >> "
        "def\n/CMapName /Adobe-Identity-UCS def\n/CMapType 2 def\n"
        "1 begincodespacerange\n<0000> <FFFF>\nendcodespacerange\n"
        f"{body}\n"
        "endcmap\nCMapName currentdict /CMap defineresource pop\nend\nend"
    ).encode("latin-1", "replace")


def find_font() -> Path | None:
    for p in FONT_PATHS:
        f = Path(p)
        if f.exists():
            return f
    return None


# ═══════════════════════════ инструменты ═══════════════════════════
def build(ws: Workspace, store: Store | None = None,
          run_id_getter: Callable[[], int] | None = None) -> list[Tool]:

    def _prep(path: str, ext: str) -> Path:
        p = ws.resolve(path)
        if p.suffix.lower() != ext:
            p = p.with_suffix(ext)
        p.parent.mkdir(parents=True, exist_ok=True)
        return p

    def _note(p: Path, kind: str, extra: str = "") -> str:
        size = p.stat().st_size
        if store is not None and run_id_getter is not None:
            store.upsert_entity("file", ws.relative(p),
                                {"kind": kind, "bytes": size},
                                run_id=run_id_getter())
        return (f"Создан {kind}: {ws.relative(p)} ({size:,} байт)"
                + (f"\n{extra}" if extra else ""))

    def make_docx(path: str, content: str, title: str = "") -> str:
        blocks = parse_markdown(content)
        if not blocks and not title:
            raise ToolError("Пустой документ создавать нечего")
        p = _prep(path, ".docx")
        if not _docx_lib(p, blocks, title):
            _docx_plain(p, blocks, title)
        tables = sum(1 for b in blocks if b.kind == "table")
        return _note(p, "документ Word",
                     f"блоков: {len(blocks)}, таблиц: {tables}")

    def make_xlsx(path: str, content: str, sheet: str = "") -> str:
        blocks = parse_markdown(content)
        sheets = _sheets_from_blocks(blocks)
        if not sheets:
            # Голый CSV тоже принимаем: модели часто отдают именно его.
            rows = [ln.split(";" if ";" in ln else ",")
                    for ln in content.strip().splitlines() if ln.strip()]
            if not rows:
                raise ToolError(
                    "Не найдено ни одной таблицы. Дайте таблицу Markdown "
                    "(| а | б | и строка-разделитель) или CSV.")
            width = max(len(r) for r in rows)
            sheets = [(_sheet_name(sheet or "Лист1", []),
                       [[c.strip() for c in r] + [""] * (width - len(r))
                        for r in rows])]
        elif sheet and len(sheets) == 1:
            sheets = [(_sheet_name(sheet, []), sheets[0][1])]
        p = _prep(path, ".xlsx")
        if not _xlsx_lib(p, sheets):
            _xlsx_plain(p, sheets)
        detail = ", ".join(f"{n}: {len(r)} строк" for n, r in sheets)
        return _note(p, "таблица Excel", f"листов {len(sheets)} — {detail}")

    def make_pptx(path: str, content: str, title: str = "") -> str:
        blocks = parse_markdown(content)
        slides = _slides_from_blocks(blocks)
        if title:
            slides.insert(0, {"title": title, "bullets": [], "rows": []})
        if not slides:
            raise ToolError(
                "Нет ни одного слайда. Новый слайд начинается с заголовка "
                "(# или ##) либо с разрыва ---")
        p = _prep(path, ".pptx")
        if not _pptx_lib(p, slides):
            _pptx_plain(p, slides)
        return _note(p, "презентация", f"слайдов: {len(slides)}")

    def make_pdf(path: str, content: str, title: str = "") -> str:
        blocks = parse_markdown(content)
        if not blocks and not title:
            raise ToolError("Пустой документ создавать нечего")
        font = find_font()
        if font is None:
            raise ToolError(
                "Не найден шрифт TrueType — без него кириллица в PDF "
                "превратится в пустые квадраты. Поставьте fonts-dejavu-core "
                "или сохраните документ в docx.")
        p = _prep(path, ".pdf")
        _make_pdf(p, blocks, title, font)
        return _note(p, "PDF", f"шрифт: {font.name}")

    return [
        Tool("make_docx",
             "Создать документ Word из Markdown. Понимает заголовки (#), "
             "абзацы, списки (-), таблицы (| а | б |) и разрыв страницы (---).",
             {"type": "object",
              "properties": {
                  "path": {"type": "string", "description": "Куда сохранить"},
                  "content": {"type": "string", "description": "Текст в Markdown"},
                  "title": {"type": "string", "description": "Заголовок сверху"}},
              "required": ["path", "content"]},
             make_docx),
        Tool("make_xlsx",
             "Создать таблицу Excel. На вход таблицы Markdown или CSV. "
             "Каждая таблица становится отдельным листом, имя берётся из "
             "заголовка над ней. Числа записываются числами, а не текстом.",
             {"type": "object",
              "properties": {
                  "path": {"type": "string"},
                  "content": {"type": "string",
                              "description": "Таблицы Markdown или CSV"},
                  "sheet": {"type": "string", "description": "Имя листа"}},
              "required": ["path", "content"]},
             make_xlsx),
        Tool("make_pptx",
             "Создать презентацию PowerPoint. Новый слайд начинается с "
             "заголовка (# или ##) либо с разрыва ---. Пункты списка и "
             "абзацы становятся содержимым слайда, таблица — таблицей.",
             {"type": "object",
              "properties": {
                  "path": {"type": "string"},
                  "content": {"type": "string"},
                  "title": {"type": "string",
                            "description": "Титульный слайд"}},
              "required": ["path", "content"]},
             make_pptx),
        Tool("make_pdf",
             "Создать PDF из Markdown. Кириллица поддерживается: шрифт "
             "встраивается в файл, текст можно искать и копировать.",
             {"type": "object",
              "properties": {
                  "path": {"type": "string"},
                  "content": {"type": "string"},
                  "title": {"type": "string"}},
              "required": ["path", "content"]},
             make_pdf),
    ]
