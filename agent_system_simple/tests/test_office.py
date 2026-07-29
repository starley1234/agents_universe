"""Тесты создания документов и связи с внешним миром.

Два правила проверки, без которых тесты были бы обманом:

  1. Файл проверяется ЧТЕНИЕМ ОБРАТНО, а не фактом создания. «Файл
     существует» ничего не значит: пустой ZIP тоже существует.
  2. Отправка проверяется на НАСТОЯЩЕМ SMTP-сервере в этом же процессе.
     Заглушка подтвердила бы только то, что мы вызвали свою же функцию;
     здесь письмо реально уходит по сокету и его содержимое сверяется.

Отдельное внимание — отказам. Для отправки наружу ложное «отправлено»
опаснее любой другой ошибки, поэтому белый список проверяется прямо.
"""
from __future__ import annotations

import asyncio
import email
import email.policy
import socket
import sys
import tempfile
import threading
import zipfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from agent.skills import comms as C                      # noqa: E402
from agent.skills import makedocs as M                   # noqa: E402
from agent.tools.base import ToolError, Workspace        # noqa: E402

PASS, FAIL = 0, 0


def check(name: str, cond: bool, detail: str = "") -> None:
    global PASS, FAIL
    if cond:
        PASS += 1
        print(f"  ok   {name}")
    else:
        FAIL += 1
        print(f"  FAIL {name}" + (f" — {detail}" if detail else ""))


def section(t: str) -> None:
    print(f"\n{t}\n" + "─" * len(t))


def free_port() -> int:
    s = socket.socket()
    s.bind(("127.0.0.1", 0))
    p = s.getsockname()[1]
    s.close()
    return p


SRC = """# Отчёт о приёмке
Изделие принято с замечаниями.

## Результаты измерений
- зазор в норме
- покрытие с царапиной

| Параметр | Норма | Факт |
|---|---|---|
| Зазор, мм | 0.5 | 0.48 |
| Масса, кг | 12 | 12.3 |

---

## Выводы
Требуется доработка покрытия.
"""


# ══════════════════════ разбор Markdown ═════════════════════════════
def test_parse() -> None:
    section("Разбор Markdown")
    b = M.parse_markdown(SRC)
    kinds = [x.kind for x in b]
    check("заголовки найдены", kinds.count("heading") == 3, str(kinds))
    check("таблица найдена", kinds.count("table") == 1, str(kinds))
    check("список найден", kinds.count("list") == 1, str(kinds))
    check("разрыв найден", kinds.count("break") == 1, str(kinds))

    tbl = [x for x in b if x.kind == "table"][0]
    check("в таблице 3 строки", len(tbl.rows) == 3, str(len(tbl.rows)))
    check("строка-разделитель не попала в данные",
          all("---" not in c for r in tbl.rows for c in r), str(tbl.rows))
    check("ячейки разобраны верно", tbl.rows[1] == ["Зазор, мм", "0.5", "0.48"],
          str(tbl.rows[1]))

    # разметка выделения снимается, ссылка превращается в текст
    b2 = M.parse_markdown("**жирный** и *косой* и [ссылка](http://x) и `код`")
    check("разметка выделения снята",
          b2[0].text == "жирный и косой и ссылка и код", b2[0].text)

    # кривая таблица не должна ронять разбор
    b3 = M.parse_markdown("| а | б |\n|---|\n| 1 |")
    t3 = [x for x in b3 if x.kind == "table"]
    check("неровная таблица выровнена по ширине",
          bool(t3) and len({len(r) for r in t3[0].rows}) == 1,
          str(t3[0].rows if t3 else "нет"))


# ═══════════════════════ создание файлов ════════════════════════════
def _tools(td: str):
    return {x.name: x for x in M.build(Workspace(td))}


def test_docx() -> None:
    section("Word: файл читается обратно")
    with tempfile.TemporaryDirectory() as td:
        t = _tools(td)
        out = t["make_docx"].fn(path="rep.docx", content=SRC, title="Приёмка")
        p = Path(td) / "rep.docx"
        check("файл создан и не пуст", p.exists() and p.stat().st_size > 1000,
              out)
        check("это настоящий ZIP-контейнер OOXML",
              zipfile.is_zipfile(p) and "word/document.xml" in
              zipfile.ZipFile(p).namelist())
        try:
            import docx                                  # type: ignore
            d = docx.Document(str(p))
            texts = [x.text for x in d.paragraphs]
            check("заголовок на месте", "Приёмка" in texts, str(texts[:3]))
            check("абзац на месте",
                  any("замечаниями" in x for x in texts))
            check("таблица сохранена как таблица", len(d.tables) == 1,
                  str(len(d.tables)))
            check("данные таблицы верны",
                  [c.text for c in d.tables[0].rows[1].cells]
                  == ["Зазор, мм", "0.5", "0.48"],
                  str([c.text for c in d.tables[0].rows[1].cells]))
        except ImportError:
            check("python-docx недоступен — чтение пропущено", True)

        # расширение подставляется само
        t["make_docx"].fn(path="без_расширения", content="текст")
        check("расширение .docx добавлено",
              (Path(td) / "без_расширения.docx").exists())
        try:
            t["make_docx"].fn(path="x.docx", content="   ")
            check("пустой документ отвергнут", False, "создан")
        except ToolError:
            check("пустой документ отвергнут", True)


def test_xlsx() -> None:
    section("Excel: числа остаются числами")
    with tempfile.TemporaryDirectory() as td:
        t = _tools(td)
        out = t["make_xlsx"].fn(path="d.xlsx", content=SRC)
        p = Path(td) / "d.xlsx"
        check("файл создан", p.exists() and zipfile.is_zipfile(p), out)
        try:
            import openpyxl                              # type: ignore
            wb = openpyxl.load_workbook(str(p))
            ws = wb.worksheets[0]
            rows = [[c.value for c in r] for r in ws.iter_rows()]
            check("имя листа взято из заголовка",
                  ws.title == "Результаты измерений", ws.title)
            check("шапка на месте", rows[0] == ["Параметр", "Норма", "Факт"],
                  str(rows[0]))
            # главное: 0.48 должно быть числом, иначе сумма не посчитается
            check("числа записаны числами, а не текстом",
                  isinstance(rows[1][2], float) and abs(rows[1][2] - 0.48) < 1e-9,
                  f"{rows[1][2]!r} ({type(rows[1][2]).__name__})")
            check("текстовая ячейка осталась текстом",
                  isinstance(rows[1][0], str), str(type(rows[1][0])))
        except ImportError:
            check("openpyxl недоступен — чтение пропущено", True)

        # CSV тоже принимается: модели часто отдают именно его
        t["make_xlsx"].fn(path="csv.xlsx", content="а;б\n1;2\n3;4")
        check("CSV принят", (Path(td) / "csv.xlsx").exists())
        try:
            t["make_xlsx"].fn(path="bad.xlsx", content="   ")
            check("пустая таблица отвергнута", False, "создана")
        except ToolError:
            check("пустая таблица отвергнута", True)

        # запрещённые символы в имени листа
        long_name = "# Отчёт: план/факт [2026]\n\n| а | б |\n|---|---|\n| 1 | 2 |"
        t["make_xlsx"].fn(path="n.xlsx", content=long_name)
        try:
            import openpyxl                              # type: ignore
            title = openpyxl.load_workbook(
                str(Path(td) / "n.xlsx")).worksheets[0].title
            check("запрещённые символы в имени листа заменены",
                  not set(title) & set(':\\/?*[]') and len(title) <= 31, title)
        except ImportError:
            pass


def test_pptx() -> None:
    section("PowerPoint: слайды разделены верно")
    with tempfile.TemporaryDirectory() as td:
        t = _tools(td)
        out = t["make_pptx"].fn(path="p.pptx", content=SRC, title="Приёмка")
        p = Path(td) / "p.pptx"
        check("файл создан", p.exists() and zipfile.is_zipfile(p), out)
        names = zipfile.ZipFile(p).namelist()
        check("слайды лежат внутри контейнера",
              sum(1 for n in names if n.startswith("ppt/slides/slide")) >= 3,
              str([n for n in names if "slide" in n]))
        try:
            from pptx import Presentation                # type: ignore
            pr = Presentation(str(p))
            slides = list(pr.slides)
            check("титульный слайд добавлен отдельно", len(slides) == 4,
                  str(len(slides)))
            titles = [s.shapes.title.text if s.shapes.title else ""
                      for s in slides]
            check("заголовки слайдов на месте",
                  "Выводы" in titles and "Приёмка" in titles, str(titles))
            check("таблица попала на слайд",
                  any(sh.has_table for s in slides for sh in s.shapes))
        except ImportError:
            check("python-pptx недоступен — чтение пропущено", True)

        try:
            t["make_pptx"].fn(path="e.pptx", content="   ")
            check("презентация без слайдов отвергнута", False, "создана")
        except ToolError:
            check("презентация без слайдов отвергнута", True)


def test_pdf() -> None:
    section("PDF: кириллица и извлекаемый текст")
    font = M.find_font()
    if font is None:
        check("шрифт TrueType не найден — тесты PDF пропущены", True)
        return
    with tempfile.TemporaryDirectory() as td:
        t = _tools(td)
        out = t["make_pdf"].fn(path="r.pdf", content=SRC, title="Приёмка")
        p = Path(td) / "r.pdf"
        raw = p.read_bytes()
        check("файл создан", p.exists() and len(raw) > 5000, out)
        check("это PDF", raw.startswith(b"%PDF-"), raw[:8].decode("latin-1"))
        check("файл завершён корректно", raw.rstrip().endswith(b"%%EOF"))
        check("шрифт встроен в файл", b"/FontFile2" in raw)
        check("карта ToUnicode есть (иначе текст не найти)",
              b"/ToUnicode" in raw and b"beginbfchar" in raw)
        check("страниц больше одной (разрыв сработал)",
              raw.count(b"/Type /Page\n") >= 2 or raw.count(b"/Type /Page ") >= 2
              or b"/Count 2" in raw, "разрыв не создал страницу")
        try:
            from pypdf import PdfReader                  # type: ignore
            r = PdfReader(str(p))
            text = "\n".join(pg.extract_text() for pg in r.pages)
            check("кириллица извлекается сторонней библиотекой",
                  "Приёмка" in text and "замечаниями" in text, text[:120])
            check("таблица попала в текст", "0.48" in text, text[:200])
            check("список отмечен маркером", "•" in text)
            check("разрыв создал вторую страницу", len(r.pages) == 2,
                  str(len(r.pages)))
        except ImportError:
            check("pypdf недоступен — извлечение пропущено", True)


def test_no_libraries() -> None:
    section("Документы создаются БЕЗ сторонних библиотек")
    # Прячем библиотеки и требуем, чтобы всё работало на stdlib.
    import builtins
    real = builtins.__import__
    hidden = {"docx", "openpyxl", "pptx"}

    def fake(name, *a, **k):
        if name.split(".")[0] in hidden:
            raise ImportError(f"скрыто тестом: {name}")
        return real(name, *a, **k)

    saved = {m: sys.modules.pop(m) for m in list(sys.modules)
             if m.split(".")[0] in hidden}
    builtins.__import__ = fake
    try:
        with tempfile.TemporaryDirectory() as td:
            t = _tools(td)
            t["make_docx"].fn(path="a.docx", content=SRC, title="Резерв")
            t["make_xlsx"].fn(path="a.xlsx", content=SRC)
            t["make_pptx"].fn(path="a.pptx", content=SRC, title="Резерв")
            for f in ("a.docx", "a.xlsx", "a.pptx"):
                p = Path(td) / f
                check(f"{f} создан на stdlib",
                      p.exists() and zipfile.is_zipfile(p) and
                      p.stat().st_size > 500, str(p.stat().st_size
                                                  if p.exists() else 0))
            keep = {f: (Path(td) / f).read_bytes()
                    for f in ("a.docx", "a.xlsx", "a.pptx")}
    finally:
        builtins.__import__ = real
        sys.modules.update(saved)

    # А теперь читаем НАСТОЯЩИМИ библиотеками: файл должен быть валидным
    with tempfile.TemporaryDirectory() as td:
        for f, data in keep.items():
            (Path(td) / f).write_bytes(data)
        try:
            import docx                                  # type: ignore
            d = docx.Document(str(Path(td) / "a.docx"))
            check("docx со stdlib читается python-docx",
                  any("замечаниями" in p.text for p in d.paragraphs))
            check("таблица из stdlib-docx читается", len(d.tables) == 1,
                  str(len(d.tables)))
        except ImportError:
            pass
        try:
            import openpyxl                              # type: ignore
            ws = openpyxl.load_workbook(str(Path(td) / "a.xlsx")).worksheets[0]
            rows = [[c.value for c in r] for r in ws.iter_rows()]
            check("xlsx со stdlib читается openpyxl",
                  rows[0] == ["Параметр", "Норма", "Факт"], str(rows[:1]))
            check("числа из stdlib-xlsx остались числами",
                  isinstance(rows[1][2], float), str(type(rows[1][2])))
        except ImportError:
            pass
        try:
            from pptx import Presentation                # type: ignore
            pr = Presentation(str(Path(td) / "a.pptx"))
            check("pptx со stdlib читается python-pptx",
                  len(list(pr.slides)) >= 3, str(len(list(pr.slides))))
        except ImportError:
            pass


# ═══════════════════════ связь: белый список ════════════════════════
def test_whitelist() -> None:
    section("Отправка: белый список")
    check("пустой список запрещает всё",
          not C.allowed(C.CommsConfig(), "a@b.ru"))
    cfg = C.CommsConfig(allow_to=["boss@firm.ru", "*@partner.com", "-100123"])
    check("точный адрес разрешён", C.allowed(cfg, "boss@firm.ru"))
    check("регистр не важен", C.allowed(cfg, "BOSS@Firm.RU"))
    check("домен по маске разрешён", C.allowed(cfg, "kto@partner.com"))
    check("чужой адрес запрещён", not C.allowed(cfg, "someone@evil.ru"))
    check("похожий домен НЕ проходит",
          not C.allowed(cfg, "kto@partner.com.evil.ru"))
    check("id чата разрешён", C.allowed(cfg, "-100123"))
    check("пустой адрес запрещён", not C.allowed(cfg, ""))


def test_draft_instead_of_send() -> None:
    section("Отправка запрещена → черновик, а не тишина")
    with tempfile.TemporaryDirectory() as td:
        ws = Workspace(td)
        # белый список пуст: не должно уйти ничего
        t = {x.name: x for x in C.build(ws, C.CommsConfig(smtp_host="х"))}
        out = t["send_email"].fn(to="kto@to.ru", subject="Тема",
                                 body="Текст письма")
        check("сказано, что НЕ отправлено", "НЕ ОТПРАВЛЕНО" in out, out[:90])
        check("названа причина", "Белый список" in out, out[:150])
        check("прямо сказано не считать это отправкой",
              "Не считай это выполненной отправкой" in out)
        drafts = list((Path(td) / "outbox").glob("*.txt"))
        check("черновик сохранён", len(drafts) == 1, str(drafts))
        body = drafts[0].read_text(encoding="utf-8")
        check("в черновике есть адресат и текст",
              "kto@to.ru" in body and "Текст письма" in body, body[:80])

        # адрес не из списка — тоже черновик
        t2 = {x.name: x for x in C.build(
            ws, C.CommsConfig(smtp_host="х", allow_to=["boss@firm.ru"]))}
        out2 = t2["send_email"].fn(to="chuzhoy@evil.ru", subject="Т",
                                   body="Б")
        check("чужой адрес не отправляется", "НЕ ОТПРАВЛЕНО" in out2, out2[:80])
        check("в отказе показан белый список", "boss@firm.ru" in out2)

        # мессенджеры под той же защитой
        out3 = t2["send_telegram"].fn(text="привет", chat="-100999")
        check("Telegram: чужой чат не отправляется",
              "НЕ ОТПРАВЛЕНО" in out3, out3[:80])
        out4 = t2["send_max"].fn(text="привет", chat="-100999")
        check("MAX: чужой чат не отправляется",
              "НЕ ОТПРАВЛЕНО" in out4, out4[:80])

        st = t2["comms_status"].fn()
        check("состояние показывает белый список", "boss@firm.ru" in st, st)


# ══════════════════ настоящая отправка по SMTP ══════════════════════
class Sink:
    """Настоящий SMTP-сервер в этом же процессе: ловит письма целиком."""

    def __init__(self) -> None:
        self.messages: list[bytes] = []

    async def handle_DATA(self, server, session, envelope):  # noqa: N802
        self.messages.append(envelope.content)
        return "250 OK"


def test_real_smtp() -> None:
    section("Почта: письмо реально уходит по сокету")
    try:
        from aiosmtpd.controller import Controller       # type: ignore
    except ImportError:
        check("aiosmtpd не установлен — проверка на живом SMTP пропущена",
              True)
        return

    sink = Sink()
    port = free_port()
    ctrl = Controller(sink, hostname="127.0.0.1", port=port)
    ctrl.start()
    try:
        with tempfile.TemporaryDirectory() as td:
            ws = Workspace(td)
            (Path(td) / "otchet.txt").write_text("вложение", encoding="utf-8")
            cfg = C.CommsConfig(smtp_host="127.0.0.1", smtp_port=port,
                                smtp_user="agent@localhost",
                                smtp_from="agent@localhost",
                                allow_to=["boss@firm.ru"])
            import os
            os.environ["AGENT_SMTP_PASS"] = "x"
            # локальный сервер без TLS: подменяем starttls на пустышку
            import smtplib
            real_starttls = smtplib.SMTP.starttls
            smtplib.SMTP.starttls = lambda self, *a, **k: (220, b"ok")
            real_login = smtplib.SMTP.login
            smtplib.SMTP.login = lambda self, *a, **k: (235, b"ok")
            try:
                t = {x.name: x for x in C.build(ws, cfg)}
                out = t["send_email"].fn(
                    to="boss@firm.ru", subject="Отчёт готов",
                    body="Тело письма с кириллицей",
                    attachments="otchet.txt")

                check("инструмент отчитался об отправке",
                      out.startswith("Письмо отправлено"), out)
                check("сервер получил ровно одно письмо",
                      len(sink.messages) == 1, str(len(sink.messages)))
                msg = email.message_from_bytes(sink.messages[0],
                                               policy=email.policy.default)
                check("адресат верный", msg["To"] == "boss@firm.ru", msg["To"])
                check("тема на месте", msg["Subject"] == "Отчёт готов",
                      str(msg["Subject"]))
                body = msg.get_body(preferencelist=("plain",)).get_content()
                check("кириллица в теле не побилась",
                      "Тело письма с кириллицей" in body, body[:60])
                names = [p.get_filename() for p in msg.iter_attachments()]
                check("вложение приложено", names == ["otchet.txt"], str(names))

                # Отказ оператора: письмо не должно уйти. Проверяем ПО
                # СЕРВЕРУ, а не по тексту ответа — текст можно и соврать.
                sink.messages.clear()
                t2 = {x.name: x for x in C.build(ws, cfg,
                                                 confirm=lambda c, r: False)}
                try:
                    out2 = t2["send_email"].fn(to="boss@firm.ru",
                                               subject="Т", body="Б")
                except ToolError as exc:
                    out2 = f"ОШИБКА: {exc}"
                check("отказ оператора: письмо не ушло на сервер",
                      not sink.messages,
                      f"сервер получил {len(sink.messages)} писем!")
                check("отказ оператора: сказано, что не отправлено",
                      "НЕ ОТПРАВЛЕНО" in out2, out2[:90])
            finally:
                smtplib.SMTP.starttls = real_starttls
                smtplib.SMTP.login = real_login
    finally:
        ctrl.stop()


def test_split() -> None:
    section("Длинные сообщения режутся по границе строки")
    text = "\n".join(f"строка номер {i}" for i in range(1, 400))
    parts = C._split_text(text, 4000)
    check("сообщение разрезано", len(parts) > 1, str(len(parts)))
    check("каждая часть в пределах лимита",
          all(len(p) <= 4000 for p in parts),
          str([len(p) for p in parts]))
    check("ничего не потеряно",
          "".join(parts).replace("\n", "") == text.replace("\n", ""))
    check("резали по строкам, а не посреди слова",
          all(not p.startswith("омер") for p in parts))
    # одна строка длиннее предела тоже должна уместиться
    long_one = "x" * 9000
    parts2 = C._split_text(long_one, 4000)
    check("сверхдлинная строка разрезана",
          all(len(p) <= 4000 for p in parts2) and
          "".join(parts2) == long_one, str([len(p) for p in parts2]))
    check("короткий текст не режется", C._split_text("привет", 4000)
          == ["привет"])


def main() -> int:
    print("=" * 60)
    print("ТЕСТЫ: создание документов и связь с внешним миром")
    print("=" * 60)
    test_parse()
    test_docx()
    test_xlsx()
    test_pptx()
    test_pdf()
    test_no_libraries()
    test_whitelist()
    test_draft_instead_of_send()
    test_real_smtp()
    test_split()
    print("\n" + "=" * 60)
    print(f"пройдено: {PASS} · провалено: {FAIL}")
    print("=" * 60)
    return 1 if FAIL else 0


if __name__ == "__main__":
    raise SystemExit(main())
